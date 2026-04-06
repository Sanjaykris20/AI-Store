import os
import sys
import subprocess
import glob

def install(package):
    subprocess.check_call([sys.executable, '-m', 'pip', 'install', package])

try:
    import docx
except ImportError:
    install('python-docx')
    import docx

try:
    import pptx
except ImportError:
    install('python-pptx')
    import pptx

try:
    import pandas as pd
except ImportError:
    install('pandas')
    import pandas as pd

try:
    import openpyxl
except ImportError:
    install('openpyxl')

try:
    import xlrd
except ImportError:
    install('xlrd')

def read_docx(path):
    try:
        doc = docx.Document(path)
        return '\n'.join([p.text for p in doc.paragraphs if p.text.strip()])
    except Exception as e:
        return str(e)

def read_pptx(path):
    try:
        prs = pptx.Presentation(path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text.append(shape.text.strip())
        return '\n'.join([t for t in text if t])
    except Exception as e:
        return str(e)

def read_xlsx(path):
    try:
        df = pd.read_excel(path, sheet_name=None)
        res = []
        for sheet, data in df.items():
            res.append(f'Sheet: {sheet}')
            res.append(data.head(20).to_string())
        return '\n'.join(res)
    except Exception as e:
        return str(e)

def read_xls(path):
    try:
        df = pd.read_excel(path, sheet_name=None, engine='xlrd')
        res = []
        for sheet, data in df.items():
            res.append(f'Sheet: {sheet}')
            res.append(data.head(20).to_string())
        return '\n'.join(res)
    except Exception as e:
        return str(e)

directory = r'c:/Users/HP/OneDrive/Documents/Infosys Intern'

extensions = {
    '*.docx': read_docx,
    '*.pptx': read_pptx,
    '*.xlsx': read_xlsx,
    '*.xls': read_xls
}

output_file = os.path.join(directory, 'extracted_docs_summary.txt')

with open(output_file, 'w', encoding='utf-8') as f:
    for ext, func in extensions.items():
        for filepath in glob.glob(os.path.join(directory, ext)):
            filename = os.path.basename(filepath)
            f.write(f'\n{'='*50}\nFILE: {filename}\n{'='*50}\n')
            content = func(filepath)
            if content:
                f.write(content[:15000])
                if len(content) > 15000:
                    f.write('\n... [TRUNCATED] ...')

print('Done extracting into extracted_docs_summary.txt')
